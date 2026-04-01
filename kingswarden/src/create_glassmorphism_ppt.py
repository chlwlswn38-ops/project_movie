from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# 디렉토리 설정
image_dir = 'kingswarden/images/ppt_visuals/'
output_path = 'kingswarden/integrated_business_report.pptx'

# 나눔고딕 폰트명 (시스템에 따라 "NanumGothic" 또는 "Nanum Gothic")
FONT_NAME = 'NanumGothic'

# 색상 설정 (Glassmorphism)
BG_COLOR = RGBColor(15, 15, 45) # #0F0F2D
GLASS_FILL = RGBColor(255, 255, 255) # White
GLASS_TRANS = 0.15 # 15% opacity
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_SOFT = RGBColor(224, 224, 240) # #E0E0F0
ACCENT_CYAN = RGBColor(103, 232, 249) # #67E8F9

def add_glass_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = GLASS_FILL
    # transparency is not directly available in standard python-pptx for fill without XML manipulation
    # so we will use a light gray/white proxy or leave as solid if necessary, 
    # but I'll try to set transparency via XML if I can.
    # For now, let's stick to a very light gray that looks semi-transparent on dark bg.
    fill.fore_color.rgb = RGBColor(35, 35, 65) # Darker tint to simulate transparency
    
    line = shape.line
    line.color.rgb = TEXT_WHITE
    line.width = Pt(1)
    
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape

def set_font(run, size, color=TEXT_WHITE, bold=False):
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold

def add_title_and_body(slide, title_text, body_text="", is_title_slide=False):
    # Set Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Title
    if is_title_slide:
        left, top, width, height = Inches(1), Inches(3), Inches(8), Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_text
        set_font(run, 44, bold=True)
    else:
        left, top, width, height = Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        set_font(run, 32, color=ACCENT_CYAN, bold=True)

    # Body
    if body_text:
        left, top, width, height = Inches(0.5), Inches(1.2), Inches(4.5), Inches(5)
        add_glass_card(slide, left, top, width, height)
        
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in body_text.split('\n'):
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            set_font(run, 14, color=TEXT_SOFT)

def add_image_to_slide(slide, img_path, left=Inches(5.2), top=Inches(1.5), width=Inches(4.5)):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width=width)

# 시작
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 1. Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "천만 영화 달성을 위한\n종합 비즈니스 및 텍스트 인텔리전스 리포트", is_title_slide=True)

# 2. Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "프로젝트 개요: 7대 핵심 영화 현황", 
                  "• 분석 대상: 명량, 기생충, 사도, 왕과 사는 남자, 헤어질 결심, 올빼미, 남산의 부장들\n\n"
                  "• 목적: 최근 천만 영화 '왕과 사는 남자'의 흥행 메커니즘 분석 및 투자자 관점의 인사이트 도출\n\n"
                  "• 핵심 목표: 개봉 초기 모멘텀 파악 및 BEP 달성 시나리오 가시화")

# 3. Data Pipeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "6단계 데이터 전처리 파이프라인", 
                  "1. Merge: 7대 핵심 영화 데이터 통합 (934개 문서)\n"
                  "2. Clean: 결측치 및 중복 제거, 순수 텍스트 추출\n"
                  "3. Tokenization: pecab 기반 형태소 분석 및 명사 추출\n"
                  "4. Limit: 문서당 1,000자 상한 제한 적용\n"
                  "5. Stopwords: 45종 무의미 범용어 필터링\n"
                  "6. TF-IDF Matrix: 빈도 및 희소성 기반 수학적 벡터화")

# 4. Methodology
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "분석 기법 및 방법론 요약", 
                  "• TF-IDF & SVD: 1000차원 행렬을 2차원으로 압축하여 영화별 군집 지형 시각화\n\n"
                  "• K-Means 군집화: 유사 주제를 공유하는 문서를 4개 그룹(N=4)으로 수학적 분리\n\n"
                  "• 비즈니스 지표: 손익분기점(BEP) 달성률 및 주차별 관객 추이 곡선 분석")

# 5. KPI Guide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "지표 가이드: 분모와 분자의 이해", 
                  "1. BEP 달성률: 목표 관객 수(분모) 대비 실제 관객 수(분자)\n\n"
                  "2. TF-IDF 중요도: 전체 영화의 흔한 단어(분모) 대비 특정 영화의 특출난 단어(분자)\n\n"
                  "3. 군집 비율: 전체 리뷰 표본(분모) 대비 특정 반응 군집(분자)")

# 6. Chart 1: BEP 달성률
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 01: 영화별 BEP 달성률 비교", "• 왕과 사는 남자는 제작비 대비 높은 BEP에도 불구하고 초단기 돌파 성공\n\n• 명량, 기생충과 같은 초우량 흥행 패턴 공유")
add_image_to_slide(slide, image_dir + 'vis_01_bep_attainment.png')

# 7. Chart 2: 관객 추이
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 02: 주차별 누적 관객 추이", "• 개봉 3주차 가파른 '엘리베이터 곡선' 완성\n\n• 오프닝 스코어 폭발 후 안정적 우상향 유지\n\n• 롱테일 형태의 천만 관객 도달 궤적")
add_image_to_slide(slide, image_dir + 'vis_02_audience_trajectory.png')

# 8. Chart 3: TF-IDF
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 03: 핵심 키워드 (TF-IDF)", "• 연기, 배우, 역사, 연출 등 본질적 가치 집중\n\n• '침술사', '비극' 등 영화별 독창적 키워드 추출\n\n• 대중의 압도적 찬사가 수치로 증명됨")
add_image_to_slide(slide, image_dir + 'vis_03_tfidf_keywords.png')

# 9. Chart 4: K-Means
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 04: K-Means 군집 분포", "• C3(대중/메가흥행) 그룹의 압도적 비중(63%)\n\n• 특정 팬덤(Niche)이 아닌 범용 대중 장악\n\n• 연기력과 몰입감에 대한 통합적 극찬")
add_image_to_slide(slide, image_dir + 'vis_04_cluster_dist.png')

# 10. Chart 5: SVD
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 05: SVD 2차원 텍스트 지형", "• 왕과 사는 남자의 피드백이 C3 공간에 강력 응집\n\n• 타 영화들과 섞이면서도 고유의 흥행 응집도 보유\n\n• 호불호 논란이 희석된 가장 깨끗한 흥행 패턴")
add_image_to_slide(slide, image_dir + 'vis_05_svd_scatter.png')

# 11. Chart 6: Sentiment
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 06: 관객 리뷰 감성 지수", "• 긍정 지수 94%로 7대 영화 중 최고 수준\n\n• 부정 피드백의 최소화가 입소문 가속화의 원동력")
add_image_to_slide(slide, image_dir + 'vis_06_sentiment.png')

# 12. Chart 7: Buzz Efficiency
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 07: 마케팅 버즈 효율 분석", "• 검색량 대비 실관객 전환율 극대화\n\n• '실속 있는 흥행' 구조 증명\n\n• 효율적 마케팅 비용 집행의 결과")
add_image_to_slide(slide, image_dir + 'vis_07_buzz_efficiency.png')

# 13. Chart 8: Marketing Funnel
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 08: 비즈니스 퍼널 분석", "• 인지에서 티켓 구매까지의 이탈률 최소화\n\n• 고관여 관객층의 적극적 참여 유도 성공")
add_image_to_slide(slide, image_dir + 'vis_08_marketing_funnel.png')

# 14. Chart 9: Revenue Structure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 09: 투자 대비 수익 구조", "• 제작비 대비 추정 매출액 압도적 상회\n\n• 투자자 관점에서 가장 매력적인 수익 모델")
add_image_to_slide(slide, image_dir + 'vis_09_revenue_structure.png')

# 15. Chart 10: Engagement
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "시각화 10: 플랫폼별 참여도", "• 유튜브 댓글 및 SNS 공유 지수 최고치\n\n• 숏폼 기반의 연기 모멘트 확산 효과")
add_image_to_slide(slide, image_dir + 'vis_10_engagement_rate.png')

# 16. Action Plan
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "전략적 액션플랜 및 타임라인", "• 초기: 배우 연기 숏폼 집중 배포 (D+14)\n\n• 중기: C3 대중 키워드 홍보 전면 전환 (3주~5주)\n\n• 상시: 리스크 모니터링 및 GV 총력전")
add_image_to_slide(slide, image_dir + 'vis_11_action_timeline.png')

# 17. Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_and_body(slide, "비즈니스 발견사항 및 제언", 
                  "1. 연기력 중심성: 흥행의 핵심 방아쇠는 배우의 압도적 기량\n\n"
                  "2. 단문 리뷰 지배: 숏폼 친화적 극찬 반응의 빠른 확산력\n\n"
                  "3. 글로벌 스탠다드: 세련된 시나리오와 완성도에 대한 관객 니즈 증명")

# 저장
prs.save(output_path)
print(f"PPT 생성이 완료되었습니다: {output_path}")
